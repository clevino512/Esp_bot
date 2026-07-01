from typing import Any
import logging
import os
import csv
from pathlib import Path

import pdfplumber
from docx import Document
from bs4 import BeautifulSoup
import markdown
from pptx import Presentation
import openpyxl
import xlrd

from app.config.constants import ALLOWED_DOCUMENT_EXTENSIONS

logger = logging.getLogger(__name__)


class DocumentLoader:
    def __init__(self):
        self.supported_extensions = ALLOWED_DOCUMENT_EXTENSIONS

    def load(self, file_path: str) -> dict[str, Any]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()
        if ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {ext}")

        content = self._extract_content(file_path, ext)
        metadata = self._extract_metadata(file_path, ext)

        return {
            "content": content,
            "metadata": metadata,
            "raw_path": str(path),
        }

    def _extract_content(self, file_path: str, ext: str) -> str:
        if ext == ".pdf":
            return self._extract_pdf(file_path)
        elif ext == ".docx":
            return self._extract_docx(file_path)
        elif ext == ".txt":
            return self._extract_txt(file_path)
        elif ext == ".md":
            return self._extract_markdown(file_path)
        elif ext == ".html":
            return self._extract_html(file_path)
        elif ext == ".xlsx":
            return self._extract_xlsx(file_path)
        elif ext == ".xls":
            return self._extract_xls(file_path)
        elif ext == ".pptx":
            return self._extract_pptx(file_path)
        elif ext == ".ppt":
            return self._extract_pptx(file_path)
        elif ext == ".csv":
            return self._extract_csv(file_path)
        return ""

    def _extract_pdf(self, file_path: str) -> str:
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return "\n\n".join(text_parts)

    def _extract_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    def _extract_txt(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def _extract_markdown(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8") as f:
            md_content = f.read()
        html = markdown.markdown(md_content)
        soup = BeautifulSoup(html, "html.parser")
        return soup.get_text(separator="\n")

    def _extract_html(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()
        soup = BeautifulSoup(html_content, "html.parser")
        return soup.get_text(separator="\n")

    def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from Excel .xlsx files."""
        parts = []
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_parts = [f"[Feuille: {sheet_name}]"]
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    sheet_parts.append(row_text)
            parts.append("\n".join(sheet_parts))
        wb.close()
        return "\n\n".join(parts)

    def _extract_xls(self, file_path: str) -> str:
        """Extract text from legacy Excel .xls files."""
        parts = []
        wb = xlrd.open_workbook(file_path)
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            sheet_parts = [f"[Feuille: {sheet.name}]"]
            for row_idx in range(sheet.nrows):
                row_text = " | ".join(str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols))
                if row_text.strip():
                    sheet_parts.append(row_text)
            parts.append("\n".join(sheet_parts))
        return "\n\n".join(parts)

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint .pptx files."""
        parts = []
        prs = Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"[Diapositive {slide_idx}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text)
            parts.append("\n".join(slide_parts))
        return "\n\n".join(parts)

    def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV files."""
        parts = []
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = " | ".join(cell for cell in row)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)

    def _extract_metadata(self, file_path: str, ext: str) -> dict[str, Any]:
        path = Path(file_path)
        stat = path.stat()
        return {
            "filename": path.name,
            "extension": ext,
            "file_size": stat.st_size,
            "modified_time": stat.st_mtime,
        }

    def load_from_directory(self, dir_path: str) -> list[dict[str, Any]]:
        documents = []
        path = Path(dir_path)

        if not path.is_dir():
            raise NotADirectoryError(f"Directory not found: {dir_path}")

        for file_path in path.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_extensions:
                try:
                    doc = self.load(str(file_path))
                    documents.append(doc)
                except Exception as e:
                    logger.error(f"Failed to load {file_path}: {e}")

        logger.info(f"Loaded {len(documents)} documents from {dir_path}")
        return documents

    def load_bytes(
        self,
        file_data: bytes,
        filename: str,
    ) -> dict[str, Any]:
        import tempfile
        ext = Path(filename).suffix.lower()

        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name

        try:
            doc = self.load(tmp_path)
            doc["metadata"]["filename"] = filename
            return doc
        finally:
            os.unlink(tmp_path)
